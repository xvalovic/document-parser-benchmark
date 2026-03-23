#!/usr/bin/env python3
from __future__ import annotations

import argparse
import asyncio
import dataclasses
import hashlib
import importlib.util
import json
import mimetypes
import os
import re
import shutil
import subprocess
import sys
import time
import traceback
from abc import ABC, abstractmethod
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any, Iterable, Optional

from dotenv import load_dotenv


@dataclasses.dataclass
class ParseResult:
    parser_name: str
    input_path: str
    success: bool
    markdown: str = ""
    plain_text: str = ""
    metadata: dict[str, Any] = dataclasses.field(default_factory=dict)
    raw: dict[str, Any] = dataclasses.field(default_factory=dict)
    elapsed_seconds: float = 0.0
    skipped_reason: Optional[str] = None
    error: Optional[str] = None


@dataclasses.dataclass
class ParserContext:
    output_root: Path
    force_ocr: bool = False
    timeout_seconds: int = 1800
    azure_use_markdown: bool = True
    textract_s3_bucket: Optional[str] = None
    textract_s3_prefix: str = "parser-benchmark"
    llamaparse_tier: str = "agentic"
    llamaparse_version: str = "latest"


IMG = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp"}
PDF = {".pdf"}
PPT = {".pptx", ".ppt"}
DOC = {".docx", ".doc"}
TXT = {".md", ".markdown", ".html", ".htm", ".txt", ".xml", ".csv"}


def ensure_dir(path: Path) -> Path:
    """Create a directory path if needed and return it."""
    path.mkdir(parents=True, exist_ok=True)
    return path


def ext(path: Path) -> str:
    """Return the lowercase file extension for a path."""
    return path.suffix.lower()


def looks_like_pdf(path: Path) -> bool:
    """Return True when a path has a PDF extension."""
    return ext(path) in PDF


def looks_like_image(path: Path) -> bool:
    """Return True when a path has a supported image extension."""
    return ext(path) in IMG


def looks_like_pptx(path: Path) -> bool:
    """Return True when a path has a PowerPoint extension."""
    return ext(path) in PPT


def looks_like_docx(path: Path) -> bool:
    """Return True when a path has a Word document extension."""
    return ext(path) in DOC


def slugify(value: str) -> str:
    """Convert a string into a filesystem-safe slug."""
    return re.sub(r"[^a-zA-Z0-9._-]+", "_", value).strip("_") or "item"


def write_text(path: Path, content: str) -> None:
    """Write UTF-8 text content to disk, creating parent directories."""
    ensure_dir(path.parent)
    path.write_text(content, encoding="utf-8")


def write_json(path: Path, payload: Any) -> None:
    """Write pretty-printed JSON to disk with UTF-8 encoding."""
    ensure_dir(path.parent)
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")


def module_exists(name: str) -> bool:
    """Check whether a Python module can be imported."""
    return importlib.util.find_spec(name) is not None


def normalize_text(text: str) -> str:
    """Normalize line endings and collapse excessive whitespace."""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def markdown_to_plain_text(md: str) -> str:
    """Convert markdown-like content into normalized plain text."""
    text = md
    text = re.sub(r"```.*?```", " ", text, flags=re.DOTALL)
    text = re.sub(r"`([^`]*)`", r"\1", text)
    text = re.sub(r"!\[[^\]]*\]\([^)]*\)", " ", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"^#{1,6}\s*", "", text, flags=re.MULTILINE)
    text = re.sub(r"^>\s?", "", text, flags=re.MULTILINE)
    text = re.sub(r"[*_~]", "", text)
    text = re.sub(r"<[^>]+>", " ", text)
    return normalize_text(text)


def collect_files(inputs: Iterable[str]) -> list[Path]:
    """Resolve input files and recursively collect files from directories."""
    out: list[Path] = []
    for raw in inputs:
        p = Path(raw).expanduser().resolve()
        if p.is_file():
            out.append(p)
        elif p.is_dir():
            out.extend([c for c in p.rglob("*") if c.is_file()])
    seen = set()
    ordered: list[Path] = []
    for p in out:
        if p not in seen:
            seen.add(p)
            ordered.append(p)
    return ordered


def infer_file_type(path: Path) -> str:
    """Infer a normalized file type label from file extension or MIME type."""
    if looks_like_pdf(path):
        return "pdf"
    if looks_like_pptx(path):
        return "pptx"
    if looks_like_image(path):
        return "image"
    if looks_like_docx(path):
        return "docx"
    if ext(path) in TXT:
        return ext(path).lstrip(".")
    guessed, _ = mimetypes.guess_type(str(path))
    return guessed or "unknown"


def heading_lines(md: str) -> list[str]:
    """Extract markdown heading lines from markdown text."""
    return [l.strip() for l in md.splitlines() if re.match(r"^#{1,6}\s+", l.strip())]


def calc_markdown_features(md: str, plain_text: str) -> dict[str, Any]:
    """Compute structural and textual metrics for a parse output."""
    return {
        "markdown_chars": len(md),
        "text_chars": len(plain_text),
        "word_count": len(re.findall(r"\b\w+\b", plain_text)),
        "heading_count": len(heading_lines(md)),
        "table_like_count": len(re.findall(r"^\|.*\|$", md, flags=re.MULTILINE))
        + len(re.findall(r"<table\b", md, flags=re.I)),
        "image_ref_count": len(re.findall(r"!\[[^\]]*\]\([^)]*\)", md))
        + len(re.findall(r"<img\b|<figure\b", md, flags=re.I)),
        "sha1_plain_text": hashlib.sha1(
            plain_text.encode("utf-8", errors="ignore")
        ).hexdigest(),
    }


def simple_similarity(a: str, b: str) -> float:
    """Compute text similarity using rapidfuzz or a bag-of-words fallback."""
    try:
        from rapidfuzz import fuzz

        return float(fuzz.ratio(a, b)) / 100.0
    except Exception:
        if not a and not b:
            return 1.0
        if not a or not b:
            return 0.0
        aw = Counter(re.findall(r"\w+", a.lower()))
        bw = Counter(re.findall(r"\w+", b.lower()))
        shared = sum((aw & bw).values())
        denom = max(1, sum(aw.values()) + sum(bw.values()))
        return (2.0 * shared) / denom


def maybe_load_ground_truth(doc_path: Path):
    """Load optional text and heading ground-truth files for a document."""
    gt = None
    heads = None
    for c in [
        doc_path.with_suffix(doc_path.suffix + ".groundtruth.txt"),
        doc_path.with_suffix(".groundtruth.txt"),
    ]:
        if c.exists():
            gt = c.read_text(encoding="utf-8")
            break
    for c in [
        doc_path.with_suffix(doc_path.suffix + ".headings.json"),
        doc_path.with_suffix(".headings.json"),
    ]:
        if c.exists():
            heads = json.loads(c.read_text(encoding="utf-8"))
            break
    return gt, heads


def evaluate_against_ground_truth(result: ParseResult, gt_text, gt_headings):
    """Evaluate a parse result against available ground-truth signals."""
    out: dict[str, Any] = {}
    if gt_text is not None:
        out["text_similarity_to_ground_truth"] = round(
            simple_similarity(
                normalize_text(result.plain_text), normalize_text(gt_text)
            ),
            4,
        )
    if gt_headings is not None:
        parsed = {
            re.sub(r"^#{1,6}\s*", "", h).strip().lower()
            for h in heading_lines(result.markdown)
        }
        truth = {str(h).strip().lower() for h in gt_headings}
        matched = len(parsed & truth)
        out["heading_recall"] = round(matched / max(1, len(truth)), 4)
        out["heading_precision"] = round(matched / max(1, len(parsed)), 4)
    return out


def run_subprocess(cmd: list[str], timeout: int):
    """Run a subprocess command with captured output and timeout."""
    return subprocess.run(
        cmd, text=True, capture_output=True, timeout=timeout, check=False
    )


class ParserAdapter(ABC):
    name = "base"

    def available(self):
        """Return whether this parser is available in the current environment."""
        return True, None

    def supports(self, path: Path) -> bool:
        """Return whether this parser can handle the given file path."""
        return True

    @abstractmethod
    def parse(self, path: Path, context: ParserContext) -> ParseResult:
        """Parse a file and return a normalized parse result."""
        raise NotImplementedError

    def skip(self, path: Path, reason: str) -> ParseResult:
        """Return a standardized skipped parse result with a reason."""
        return ParseResult(
            parser_name=self.name,
            input_path=str(path),
            success=False,
            skipped_reason=reason,
        )


class DoclingAdapter(ParserAdapter):
    name = "docling"

    def available(self):
        """Check whether the docling package is installed."""
        ok = module_exists("docling")
        return ok, None if ok else "pip install docling"

    def parse(self, path, context):
        """Parse a document using Docling and export markdown."""
        start = time.perf_counter()
        from docling.document_converter import DocumentConverter

        doc = DocumentConverter().convert(str(path)).document
        md = doc.export_to_markdown()
        return ParseResult(
            self.name,
            str(path),
            True,
            normalize_text(md),
            markdown_to_plain_text(md),
            {"format": infer_file_type(path)},
            {},
            time.perf_counter() - start,
        )


class PyMuPDF4LLMAdapter(ParserAdapter):
    name = "pymupdf4llm"

    def available(self):
        """Check whether pymupdf4llm is installed."""
        ok = module_exists("pymupdf4llm")
        return ok, None if ok else "pip install pymupdf4llm"

    def supports(self, path):
        """Limit this adapter to PDF inputs."""
        return looks_like_pdf(path)

    def parse(self, path, context):
        """Parse a PDF into markdown with pymupdf4llm."""
        start = time.perf_counter()
        import pymupdf4llm

        md = pymupdf4llm.to_markdown(
            str(path), page_chunks=False, use_ocr=True, force_ocr=context.force_ocr
        )
        return ParseResult(
            self.name,
            str(path),
            True,
            normalize_text(md),
            markdown_to_plain_text(md),
            {"format": "pdf"},
            {},
            time.perf_counter() - start,
        )


class UnstructuredAdapter(ParserAdapter):
    name = "unstructured"

    def available(self):
        """Check whether unstructured is installed."""
        ok = module_exists("unstructured")
        return ok, None if ok else "pip install 'unstructured[all-docs]'"

    def parse(self, path, context):
        """Parse a document with unstructured and map elements to markdown."""
        start = time.perf_counter()
        from unstructured.partition.auto import partition

        elements = partition(filename=str(path))
        lines = []
        preview = []
        for el in elements:
            category = getattr(el, "category", None) or el.__class__.__name__
            text = str(el).strip()
            if not text:
                continue
            if str(category).lower() in {"title", "header", "section-header"}:
                lines.append(f"# {text}")
            elif str(category).lower() in {"listitem", "list-item"}:
                lines.append(f"- {text}")
            else:
                lines.append(text)
            preview.append({"category": str(category), "text": text[:500]})
        md = "\n\n".join(lines)
        return ParseResult(
            self.name,
            str(path),
            True,
            normalize_text(md),
            markdown_to_plain_text(md),
            {"element_count": len(elements)},
            {"elements_preview": preview[:100]},
            time.perf_counter() - start,
        )


class NativePptxAdapter(ParserAdapter):
    name = "native_pptx"

    def available(self):
        """Check whether python-pptx is installed."""
        ok = module_exists("pptx")
        return ok, None if ok else "pip install python-pptx"

    def supports(self, path):
        """Limit this adapter to PowerPoint files."""
        return looks_like_pptx(path)

    def parse(self, path, context):
        """Extract slide text from PPT/PPTX files into markdown."""
        start = time.perf_counter()
        from pptx import Presentation

        prs = Presentation(str(path))
        parts = [f"# {path.name}"]
        slides = []
        for idx, slide in enumerate(prs.slides, start=1):
            title = None
            texts = []
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
                if not text:
                    continue
                if title is None:
                    title = text.splitlines()[0].strip()
                texts.append(text)
            parts.append(f"## Slide {idx}: {title or 'Untitled'}")
            parts.extend(texts)
            slides.append(
                {"slide": idx, "title": title, "text_block_count": len(texts)}
            )
        md = "\n\n".join(parts)
        return ParseResult(
            self.name,
            str(path),
            True,
            normalize_text(md),
            markdown_to_plain_text(md),
            {"slide_count": len(prs.slides)},
            {"slides": slides},
            time.perf_counter() - start,
        )


class NativeDocxAdapter(ParserAdapter):
    name = "native_docx"

    def available(self):
        """Check whether python-docx is installed."""
        ok = module_exists("docx")
        return ok, None if ok else "pip install python-docx"

    def supports(self, path):
        """Limit this adapter to Word document files."""
        return looks_like_docx(path)

    def parse(self, path, context):
        """Extract paragraphs and tables from DOC/DOCX into markdown."""
        start = time.perf_counter()
        import docx

        doc = docx.Document(str(path))
        parts = [f"# {path.name}"]
        table_count = 0
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = getattr(getattr(para, "style", None), "name", "") or ""
            if style.lower().startswith("heading"):
                m = re.search(r"(\d+)", style)
                level = min(6, int(m.group(1))) if m else 1
                parts.append(f"{'#' * level} {text}")
            else:
                parts.append(text)
        for table in doc.tables:
            table_count += 1
            rows = []
            for row in table.rows:
                rows.append(
                    "| "
                    + " | ".join(
                        cell.text.strip().replace("\n", " ") for cell in row.cells
                    )
                    + " |"
                )
            if rows:
                parts.append(f"## Table {table_count}")
                parts.extend(rows)
        md = "\n\n".join(parts)
        return ParseResult(
            self.name,
            str(path),
            True,
            normalize_text(md),
            markdown_to_plain_text(md),
            {"table_count": table_count},
            {},
            time.perf_counter() - start,
        )


class MarkerAdapter(ParserAdapter):
    name = "marker"

    def available(self):
        """Check whether marker CLI or Python package is available."""
        if shutil.which("marker_single") or module_exists("marker"):
            return True, None
        return False, "pip install 'marker-pdf[full]'"

    def parse(self, path, context):
        """Parse files with Marker via CLI when possible, otherwise Python API."""
        start = time.perf_counter()
        outdir = ensure_dir(
            context.output_root
            / "_tmp"
            / f"marker_{slugify(path.stem)}_{int(time.time() * 1000)}"
        )
        if shutil.which("marker_single"):
            cmd = [
                "marker_single",
                str(path),
                "--output_format",
                "markdown",
                "--output_dir",
                str(outdir),
            ]
            if context.force_ocr:
                cmd.append("--force_ocr")
            proc = run_subprocess(cmd, timeout=context.timeout_seconds)
            if proc.returncode != 0:
                raise RuntimeError(
                    f"marker_single failed\nSTDOUT:\n{proc.stdout}\nSTDERR:\n{proc.stderr}"
                )
            md_files = sorted(
                outdir.rglob("*.md"), key=lambda p: p.stat().st_mtime, reverse=True
            )
            if not md_files:
                raise RuntimeError(
                    f"marker_single did not produce markdown in {outdir}"
                )
            md = md_files[0].read_text(encoding="utf-8")
            return ParseResult(
                self.name,
                str(path),
                True,
                normalize_text(md),
                markdown_to_plain_text(md),
                {"used": "marker_single"},
                {"stdout": proc.stdout[-2000:], "stderr": proc.stderr[-2000:]},
                time.perf_counter() - start,
            )
        if not looks_like_pdf(path):
            return self.skip(path, "install marker_single for non-PDF coverage")
        from marker.converters.pdf import PdfConverter
        from marker.models import create_model_dict

        rendered = PdfConverter(artifact_dict=create_model_dict())(str(path))
        md = getattr(rendered, "markdown", None)
        if not md:
            raise RuntimeError("Marker Python API returned no markdown")
        return ParseResult(
            self.name,
            str(path),
            True,
            normalize_text(md),
            markdown_to_plain_text(md),
            {"used": "python_api"},
            {},
            time.perf_counter() - start,
        )


class LlamaParseAdapter(ParserAdapter):
    name = "llamaparse"

    def available(self):
        """Check whether LlamaParse SDK and API key are available."""
        if not module_exists("llama_cloud"):
            return False, "pip install llama_cloud>=1.0"
        if not os.getenv("LLAMA_CLOUD_API_KEY") and not os.getenv(
            "LLAMA_PARSE_API_KEY"
        ):
            return False, "set LLAMA_CLOUD_API_KEY"
        return True, None

    def parse(self, path, context):
        """Run the async LlamaParse flow from a synchronous entrypoint."""
        return asyncio.run(self._parse_async(path, context))

    async def _parse_async(self, path, context):
        """Upload a file to Llama Cloud and fetch markdown parse output."""
        start = time.perf_counter()
        from llama_cloud import AsyncLlamaCloud

        api_key = os.getenv("LLAMA_CLOUD_API_KEY") or os.getenv("LLAMA_PARSE_API_KEY")
        client = AsyncLlamaCloud(api_key=api_key)
        file_obj = await client.files.create(file=str(path), purpose="parse")
        parse_result = await client.parsing.parse(
            file_id=file_obj.id,
            tier=context.llamaparse_tier,
            version=context.llamaparse_version,
            expand=["markdown"],
        )

        pages = getattr(getattr(parse_result, "markdown", None), "pages", []) or []
        md_parts = [
            getattr(page, "markdown", "")
            for page in pages
            if getattr(page, "markdown", "")
        ]
        md = "\n\n".join(md_parts)

        metadata = {
            "file_id": getattr(file_obj, "id", None),
            "parse_job_id": getattr(parse_result, "id", None),
            "tier": context.llamaparse_tier,
            "version": context.llamaparse_version,
            "page_count": len(pages),
        }
        raw = {
            "status": getattr(parse_result, "status", None),
        }
        return ParseResult(
            self.name,
            str(path),
            True,
            normalize_text(md),
            markdown_to_plain_text(md),
            metadata,
            raw,
            time.perf_counter() - start,
        )


class AzureDocumentIntelligenceAdapter(ParserAdapter):
    name = "azure_document_intelligence"

    def available(self):
        """Check whether Azure Document Intelligence prerequisites are set."""
        if not module_exists("requests"):
            return False, "pip install requests"
        if not os.getenv("DOCUMENTINTELLIGENCE_ENDPOINT"):
            return False, "set DOCUMENTINTELLIGENCE_ENDPOINT"
        if not os.getenv("DOCUMENTINTELLIGENCE_API_KEY"):
            return False, "set DOCUMENTINTELLIGENCE_API_KEY"
        return True, None

    def parse(self, path, context):
        """Submit a file to Azure Document Intelligence and poll for results."""
        start = time.perf_counter()
        import requests

        endpoint = os.environ["DOCUMENTINTELLIGENCE_ENDPOINT"].rstrip("/")
        key = os.environ["DOCUMENTINTELLIGENCE_API_KEY"]
        api_version = os.getenv("DOCUMENTINTELLIGENCE_API_VERSION", "2024-11-30")
        model_id = os.getenv("DOCUMENTINTELLIGENCE_MODEL_ID", "prebuilt-layout")
        output_format = "markdown" if context.azure_use_markdown else "text"
        url = f"{endpoint}/documentintelligence/documentModels/{model_id}:analyze?api-version={api_version}&outputContentFormat={output_format}"
        content_type = mimetypes.guess_type(str(path))[0] or "application/octet-stream"
        with open(path, "rb") as f:
            resp = requests.post(
                url,
                headers={
                    "Ocp-Apim-Subscription-Key": key,
                    "Content-Type": content_type,
                },
                data=f.read(),
                timeout=120,
            )
        resp.raise_for_status()
        op_url = resp.headers.get("operation-location")
        if not op_url:
            raise RuntimeError(
                f"Azure did not return operation-location: {resp.text[:1000]}"
            )
        deadline = time.time() + context.timeout_seconds
        while True:
            poll = requests.get(
                op_url, headers={"Ocp-Apim-Subscription-Key": key}, timeout=60
            )
            poll.raise_for_status()
            body = poll.json()
            status = body.get("status")
            if status == "succeeded":
                result = body.get("analyzeResult", body)
                content = result.get("content", "")
                return ParseResult(
                    self.name,
                    str(path),
                    True,
                    normalize_text(content),
                    markdown_to_plain_text(content),
                    {
                        "api_version": api_version,
                        "model_id": model_id,
                        "content_format": result.get("contentFormat", output_format),
                    },
                    {"analyze_result": result},
                    time.perf_counter() - start,
                )
            if status in {"failed", "canceled"}:
                raise RuntimeError(json.dumps(body)[:4000])
            if time.time() > deadline:
                raise TimeoutError(f"Azure analysis timed out for {path}")
            time.sleep(2)


class TextractAdapter(ParserAdapter):
    name = "amazon_textract"

    def available(self):
        """Check whether boto3 is installed for Textract access."""
        ok = module_exists("boto3")
        return ok, None if ok else "pip install boto3"

    def supports(self, path):
        """Limit this adapter to image and PDF inputs."""
        return looks_like_image(path) or looks_like_pdf(path)

    def parse(self, path, context):
        """Parse files with Amazon Textract and return line-level markdown."""
        start = time.perf_counter()
        import boto3

        client = boto3.client("textract")
        if looks_like_pdf(path):
            bucket = context.textract_s3_bucket or os.getenv("TEXTRACT_S3_BUCKET")
            if not bucket:
                return self.skip(
                    path,
                    "Textract PDF benchmarking requires S3 input; set TEXTRACT_S3_BUCKET",
                )
            s3 = boto3.client("s3")
            key = f"{context.textract_s3_prefix.strip('/')}/{slugify(path.name)}"
            s3.upload_file(str(path), bucket, key)
            job = client.start_document_analysis(
                DocumentLocation={"S3Object": {"Bucket": bucket, "Name": key}},
                FeatureTypes=["TABLES", "FORMS"],
            )
            job_id = job["JobId"]
            deadline = time.time() + context.timeout_seconds
            pages = []
            next_token = None
            while True:
                kwargs = {"JobId": job_id}
                if next_token:
                    kwargs["NextToken"] = next_token
                result = client.get_document_analysis(**kwargs)
                status = result["JobStatus"]
                if status == "SUCCEEDED":
                    pages.append(result)
                    next_token = result.get("NextToken")
                    if next_token:
                        continue
                    break
                if status in {"FAILED", "PARTIAL_SUCCESS"}:
                    raise RuntimeError(json.dumps(result)[:4000])
                if time.time() > deadline:
                    raise TimeoutError(f"Textract timed out for {path}")
                time.sleep(2)
            blocks = [b for p in pages for b in p.get("Blocks", [])]
        else:
            result = client.analyze_document(
                Document={"Bytes": path.read_bytes()}, FeatureTypes=["TABLES", "FORMS"]
            )
            blocks = result.get("Blocks", [])
        lines = [
            b["Text"] for b in blocks if b.get("BlockType") == "LINE" and b.get("Text")
        ]
        md = "\n\n".join(lines)
        return ParseResult(
            self.name,
            str(path),
            True,
            normalize_text(md),
            normalize_text("\n".join(lines)),
            {"block_count": len(blocks), "line_count": len(lines)},
            {"blocks_preview": blocks[:200]},
            time.perf_counter() - start,
        )


def build_adapters():
    """Build the ordered list of parser adapters used in benchmarking."""
    return [
        DoclingAdapter(),
        MarkerAdapter(),
        LlamaParseAdapter(),
        AzureDocumentIntelligenceAdapter(),
        UnstructuredAdapter(),
        PyMuPDF4LLMAdapter(),
        TextractAdapter(),
        NativePptxAdapter(),
        NativeDocxAdapter(),
    ]


def safe_parse(adapter, path, context):
    """Run parser execution with availability, support, and error handling."""
    available, reason = adapter.available()
    if not available:
        return adapter.skip(path, reason or "not available")
    if not adapter.supports(path):
        return adapter.skip(path, f"not applicable to {infer_file_type(path)}")
    try:
        return adapter.parse(path, context)
    except Exception as exc:
        return ParseResult(
            parser_name=adapter.name,
            input_path=str(path),
            success=False,
            error=f"{type(exc).__name__}: {exc}",
            raw={"traceback": traceback.format_exc(limit=20)},
        )


def save_result(path, result, row, output_root):
    """Persist parse artifacts and evaluation data for one parser/file pair."""
    base = ensure_dir(
        output_root / slugify(Path(path).name) / slugify(result.parser_name)
    )
    if result.markdown:
        write_text(base / "output.md", result.markdown)
    if result.plain_text:
        write_text(base / "output.txt", result.plain_text)
    write_json(base / "result.json", dataclasses.asdict(result))
    write_json(base / "evaluation.json", row)


def rank_results(rows):
    """Compute heuristic scores and return rows sorted best-first."""
    for row in rows:
        score = 1.0 if row.get("success") else 0.0
        if row.get("text_similarity_to_ground_truth") is not None:
            score += float(row["text_similarity_to_ground_truth"]) * 3.0
        if row.get("heading_recall") is not None:
            score += float(row["heading_recall"]) * 2.0
        score += min(float(row.get("heading_count", 0)), 20.0) / 20.0
        score += min(float(row.get("table_like_count", 0)), 10.0) / 20.0
        score -= min(float(row.get("elapsed_seconds", 0.0)), 600.0) / 600.0
        row["heuristic_score"] = round(score, 4)
    return sorted(
        rows,
        key=lambda x: (x.get("heuristic_score", 0), x.get("success", False)),
        reverse=True,
    )


def classify_skip_reason(skipped_reason: Optional[str]) -> str:
    """Normalize skip reasons into coarse categories for reporting."""
    if not skipped_reason:
        return "none"
    if skipped_reason.startswith("not applicable to "):
        return "not_applicable"
    if skipped_reason.startswith("pip install ") or skipped_reason.startswith("set "):
        return "unavailable"
    return "skipped"


def aggregate_parser_scores(rows: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Aggregate parser performance metrics across a set of rows."""
    by_parser: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for row in rows:
        by_parser[row["parser_name"]].append(row)

    leaderboard = []
    for parser_name, parser_rows in by_parser.items():
        attempted = [r for r in parser_rows if not r.get("skipped_reason")]
        successful = [r for r in attempted if r.get("success")]
        not_applicable = sum(
            1
            for r in parser_rows
            if classify_skip_reason(r.get("skipped_reason")) == "not_applicable"
        )
        unavailable = sum(
            1
            for r in parser_rows
            if classify_skip_reason(r.get("skipped_reason")) == "unavailable"
        )
        skipped_other = sum(
            1
            for r in parser_rows
            if classify_skip_reason(r.get("skipped_reason")) == "skipped"
        )

        run_count = len(attempted)
        success_count = len(successful)
        avg_score = (
            round(
                sum(float(r.get("heuristic_score", 0.0)) for r in attempted)
                / run_count,
                4,
            )
            if run_count
            else None
        )
        avg_elapsed = (
            round(
                sum(float(r.get("elapsed_seconds", 0.0)) for r in successful)
                / success_count,
                3,
            )
            if success_count
            else None
        )

        leaderboard.append(
            {
                "parser_name": parser_name,
                "files_seen": len(parser_rows),
                "run_count": run_count,
                "success_count": success_count,
                "success_rate": round(success_count / run_count, 4)
                if run_count
                else None,
                "avg_heuristic_score": avg_score,
                "avg_elapsed_seconds_success": avg_elapsed,
                "not_applicable_count": not_applicable,
                "unavailable_count": unavailable,
                "skipped_other_count": skipped_other,
            }
        )

    leaderboard.sort(
        key=lambda r: (
            r["avg_heuristic_score"] is not None,
            r["avg_heuristic_score"]
            if r["avg_heuristic_score"] is not None
            else float("-inf"),
            r["success_rate"] if r["success_rate"] is not None else float("-inf"),
            -(
                r["avg_elapsed_seconds_success"]
                if r["avg_elapsed_seconds_success"] is not None
                else float("inf")
            ),
        ),
        reverse=True,
    )
    return leaderboard


def aggregate_parser_scores_by_filetype(
    rows: list[dict[str, Any]],
) -> dict[str, list[dict[str, Any]]]:
    """Aggregate parser performance separately for each detected file type."""
    by_type: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for row in rows:
        by_type[row.get("input_type", "unknown")].append(row)

    return {
        file_type: aggregate_parser_scores(type_rows)
        for file_type, type_rows in sorted(by_type.items())
    }


def build_summary_markdown(rows):
    """Build a human-readable markdown summary of benchmark results."""
    lines = ["# Parser Benchmark Summary", ""]
    grouped = {}
    for row in rows:
        grouped.setdefault(row["input_file"], []).append(row)
    for input_file, group in grouped.items():
        lines.append(f"## {Path(input_file).name}")
        for item in rank_results(group):
            status = (
                "OK"
                if item.get("success")
                else f"SKIP/FAIL: {item.get('skipped_reason') or item.get('error', 'unknown')}"
            )
            lines.append(
                f"- **{item['parser_name']}** — {status}; score={item.get('heuristic_score')}, secs={item.get('elapsed_seconds')}, words={item.get('word_count')}, headings={item.get('heading_count')}, tables={item.get('table_like_count')}"
            )
        lines.append("")

    overall = aggregate_parser_scores(rows)
    lines.append("## Overall parser leaderboard")
    if overall:
        for idx, item in enumerate(overall, start=1):
            lines.append(
                f"{idx}. **{item['parser_name']}** — avg_score={item.get('avg_heuristic_score')}, success_rate={item.get('success_rate')}, run_count={item.get('run_count')}, unavailable={item.get('unavailable_count')}, not_applicable={item.get('not_applicable_count')}"
            )
    else:
        lines.append("- No parser runs available for leaderboard.")
    lines.append("")

    by_type = aggregate_parser_scores_by_filetype(rows)
    lines.append("## Parser leaderboard by file type")
    for file_type, leaderboard in by_type.items():
        lines.append(f"### {file_type}")
        if leaderboard:
            for idx, item in enumerate(leaderboard, start=1):
                lines.append(
                    f"{idx}. **{item['parser_name']}** — avg_score={item.get('avg_heuristic_score')}, success_rate={item.get('success_rate')}, run_count={item.get('run_count')}"
                )
        else:
            lines.append("- No parser runs available for this file type.")
        lines.append("")

    return "\n".join(lines)


def parse_args():
    """Parse CLI arguments for benchmark execution."""
    ap = argparse.ArgumentParser(
        description="POC benchmark for AI Knowledge document parsers"
    )
    ap.add_argument("inputs", nargs="+")
    ap.add_argument("--output-dir", default="benchmark_output")
    ap.add_argument("--parsers", default="all")
    ap.add_argument("--force-ocr", action="store_true")
    ap.add_argument("--timeout-seconds", type=int, default=1800)
    ap.add_argument("--azure-plain-text", action="store_true")
    ap.add_argument("--textract-s3-bucket", default=os.getenv("TEXTRACT_S3_BUCKET"))
    ap.add_argument(
        "--textract-s3-prefix",
        default=os.getenv("TEXTRACT_S3_PREFIX", "parser-benchmark"),
    )
    ap.add_argument(
        "--llamaparse-tier", default=os.getenv("LLAMAPARSE_TIER", "agentic")
    )
    ap.add_argument(
        "--llamaparse-version", default=os.getenv("LLAMAPARSE_VERSION", "latest")
    )
    return ap.parse_args()


def main():
    """Run the end-to-end benchmark pipeline for input files."""
    load_dotenv(override=False)
    args = parse_args()
    files = collect_files(args.inputs)
    if not files:
        print("No input files found.", file=sys.stderr)
        return 2
    output_root = ensure_dir(Path(args.output_dir).resolve())
    context = ParserContext(
        output_root=output_root,
        force_ocr=args.force_ocr,
        timeout_seconds=args.timeout_seconds,
        azure_use_markdown=not args.azure_plain_text,
        textract_s3_bucket=args.textract_s3_bucket,
        textract_s3_prefix=args.textract_s3_prefix,
        llamaparse_tier=args.llamaparse_tier,
        llamaparse_version=args.llamaparse_version,
    )
    adapters = build_adapters()
    selected = (
        None
        if args.parsers == "all"
        else {p.strip() for p in args.parsers.split(",") if p.strip()}
    )
    if selected is not None:
        adapters = [a for a in adapters if a.name in selected]
    rows = []
    availability = []
    for adapter in adapters:
        available, reason = adapter.available()
        availability.append(
            {"parser": adapter.name, "available": available, "details": reason}
        )
    write_json(output_root / "availability.json", availability)
    for path in files:
        gt_text, gt_headings = maybe_load_ground_truth(path)
        print(f"\n=== {path} ({infer_file_type(path)}) ===")
        for adapter in adapters:
            print(f"  -> {adapter.name}")
            result = safe_parse(adapter, path, context)
            row = {
                "input_file": str(path),
                "input_type": infer_file_type(path),
                "parser_name": adapter.name,
                "success": result.success,
                "elapsed_seconds": round(result.elapsed_seconds, 3),
                "skipped_reason": result.skipped_reason,
                "error": result.error,
            }
            if result.success:
                row.update(calc_markdown_features(result.markdown, result.plain_text))
                row.update(evaluate_against_ground_truth(result, gt_text, gt_headings))
            rows.append(row)
            save_result(path, result, row, output_root)
    ranked = rank_results(rows)
    write_json(output_root / "summary.json", ranked)
    write_json(
        output_root / "parser_leaderboard_overall.json", aggregate_parser_scores(ranked)
    )
    write_json(
        output_root / "parser_leaderboard_by_filetype.json",
        aggregate_parser_scores_by_filetype(ranked),
    )
    write_text(output_root / "summary.md", build_summary_markdown(ranked))
    print(f"\nWrote benchmark outputs to: {output_root}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
